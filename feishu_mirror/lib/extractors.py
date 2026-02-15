from __future__ import annotations

import hashlib
from dataclasses import dataclass
from pathlib import Path
from typing import Callable


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".md", ".pptx", ".xlsx"}


@dataclass
class ExtractedDocument:
    text: str
    title: str
    meta: dict


def sha256_text(content: str) -> str:
    return hashlib.sha256(content.encode("utf-8", errors="ignore")).hexdigest()


def read_markdown(path: Path) -> ExtractedDocument:
    text = path.read_text(encoding="utf-8", errors="ignore")
    title = path.stem
    for line in text.splitlines():
        stripped = line.strip()
        if stripped.startswith("#"):
            title = stripped.lstrip("#").strip() or title
            break
    return ExtractedDocument(text=text, title=title, meta={"parser": "markdown"})


def read_docx(path: Path) -> ExtractedDocument:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required for .docx parsing") from exc

    doc = Document(str(path))
    lines = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    text = "\n".join(lines)
    title = lines[0][:120] if lines else path.stem
    return ExtractedDocument(text=text, title=title, meta={"parser": "docx", "paragraphs": len(lines)})


def read_pdf(path: Path) -> ExtractedDocument:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("pypdf is required for .pdf parsing") from exc

    reader = PdfReader(str(path))
    pages: list[str] = []
    for page in reader.pages:
        pages.append((page.extract_text() or "").strip())
    text = "\n\n".join(p for p in pages if p)
    title = path.stem
    metadata = reader.metadata or {}
    if getattr(metadata, "title", None):
        title = str(metadata.title)
    return ExtractedDocument(text=text, title=title, meta={"parser": "pdf", "pages": len(reader.pages)})


def read_pptx(path: Path) -> ExtractedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for .pptx parsing") from exc

    prs = Presentation(str(path))
    lines: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                content = shape.text.strip()
                if content:
                    slide_lines.append(content)
        if slide_lines:
            lines.append(f"[Slide {idx}]\n" + "\n".join(slide_lines))
    text = "\n\n".join(lines)
    title = path.stem
    if prs.core_properties and prs.core_properties.title:
        title = str(prs.core_properties.title)
    return ExtractedDocument(text=text, title=title, meta={"parser": "pptx", "slides": len(prs.slides)})


def read_xlsx(path: Path) -> ExtractedDocument:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("openpyxl is required for .xlsx parsing") from exc

    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    sections: list[str] = []
    cells_count = 0
    for ws in wb.worksheets:
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if not vals:
                continue
            cells_count += len(vals)
            rows_text.append(" | ".join(vals))
        if rows_text:
            sections.append(f"[Sheet {ws.title}]\n" + "\n".join(rows_text))
    wb.close()
    text = "\n\n".join(sections)
    return ExtractedDocument(
        text=text,
        title=path.stem,
        meta={"parser": "xlsx", "sheets": len(wb.worksheets), "cells": cells_count},
    )


def extractor_for(path: Path) -> Callable[[Path], ExtractedDocument] | None:
    mapping: dict[str, Callable[[Path], ExtractedDocument]] = {
        ".md": read_markdown,
        ".docx": read_docx,
        ".pdf": read_pdf,
        ".pptx": read_pptx,
        ".xlsx": read_xlsx,
    }
    return mapping.get(path.suffix.lower())
